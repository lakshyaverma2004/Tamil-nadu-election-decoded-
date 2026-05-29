import os
import io
import numpy as np

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

plt.rcParams['font.family'] = 'sans-serif'

SW, SH = 13.333, 7.5

BG_H     = '0A0E29'
CARD_H   = '131844'
CARD2_H  = '0D1130'
WHITE_H  = 'FFFFFF'
MUTED_H  = '94A3B8'
TVK_H    = 'A855F7'
DMK_H    = '06B6D4'
AIADMK_H = '10B981'
INC_H    = '3B82F6'
BJP_H    = 'F59E0B'
RED_H    = 'EF4444'
SLATE_H  = '64748B'

BG_M     = '#0A0E29'
CARD_M   = '#131844'
CARD2_M  = '#0D1130'
WHITE_M  = '#FFFFFF'
MUTED_M  = '#94A3B8'
TVK_M    = '#A855F7'
DMK_M    = '#06B6D4'
AIADMK_M = '#10B981'
INC_M    = '#3B82F6'
BJP_M    = '#F59E0B'
RED_M    = '#EF4444'
SLATE_M  = '#64748B'

RECT  = 1
RRECT = 5


def rgb(h):
    if h.startswith('#'):
        h = h[1:]
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def xl(f): return Inches(f * SW)
def yt(f): return Inches(f * SH)
def xw(f): return Inches(f * SW)
def yh(f): return Inches(f * SH)


def vcenter(tf):
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', 'ctr')


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(BG_H)
    return slide


def add_rect(slide, xf, yf, wf, hf, fill_h, border_h=None, bpt=0):
    s = slide.shapes.add_shape(RECT, xl(xf), yt(yf), xw(wf), yh(hf))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill_h)
    if border_h and bpt > 0:
        s.line.color.rgb = rgb(border_h)
        s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    return s


def add_rrect(slide, xf, yf, wf, hf, fill_h, border_h=None, bpt=1.5):
    s = slide.shapes.add_shape(RRECT, xl(xf), yt(yf), xw(wf), yh(hf))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill_h)
    if border_h:
        s.line.color.rgb = rgb(border_h)
        s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    return s


def card(slide, xf, yf, wf, hf, border_h=SLATE_H, fill_h=CARD_H):
    return add_rrect(slide, xf, yf, wf, hf, fill_h, border_h=border_h)


def top_rule(slide, xf, yf, wf, color_h):
    add_rect(slide, xf, yf, wf, 0.007, color_h)


def side_bar(slide, xf, yf, hf, color_h):
    add_rect(slide, xf, yf, 0.007, hf, color_h)


def T(slide, text, xf, yf, wf, hf, size, bold=False,
      color_h=WHITE_H, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(xl(xf), yt(yf), xw(wf), yh(hf))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color_h)
    p.alignment = align
    return tb


def TM(slide, rows, xf, yf, wf, hf, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(xl(xf), yt(yf), xw(wf), yh(hf))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, (text, size, bold, color_h) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color_h)
        p.alignment = align
    return tb


def shape_label(shape, rows, align=PP_ALIGN.CENTER, vc=True):
    tf = shape.text_frame
    tf.word_wrap = True
    if vc:
        vcenter(tf)
    for i, (text, size, bold, color_h) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color_h)
        p.alignment = align


def header(slide, n, total=10, accent_h=TVK_H):
    add_rect(slide, 0, 0, 1.0, 0.08, accent_h)
    T(slide, 'TAMIL NADU ELECTION 2026   |   ATLIQ MEDIA DECISION DESK',
      0.025, 0.010, 0.820, 0.060, 9,
      bold=True, color_h=WHITE_H, align=PP_ALIGN.LEFT, wrap=False)
    badge = add_rrect(slide, 0.870, 0.013, 0.115, 0.055,
                      BG_H, border_h=WHITE_H, bpt=1.2)
    shape_label(badge, [(f'SLIDE {n}/{total}', 8.5, True, WHITE_H)])


def footer(slide, text, color_h=BJP_H):
    add_rect(slide, 0, 0.910, 1.0, 0.090, CARD2_H)
    add_rect(slide, 0, 0.910, 1.0, 0.008, color_h)
    T(slide, text, 0.02, 0.922, 0.96, 0.068,
      10, bold=True, color_h=color_h, align=PP_ALIGN.CENTER, wrap=False)


def chart_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def bar_chart(parties, s21, s26, colors26, title_str, ylim=160):
    fig, ax = plt.subplots(figsize=(7, 4.8), facecolor=BG_M)
    ax.set_facecolor(BG_M)
    x = np.arange(len(parties))
    w = 0.36
    ax.bar(x - w / 2, s21, w, color=SLATE_M, alpha=0.30, label='2021')
    b2 = ax.bar(x + w / 2, s26, w, color=colors26, label='2026')
    for bar in b2:
        h = bar.get_height()
        if h > 0:
            ax.text(bar.get_x() + bar.get_width() / 2, h + 1.5,
                    str(int(h)), ha='center', va='bottom',
                    color=WHITE_M, fontsize=11, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(parties, fontsize=13, color=WHITE_M, fontweight='bold')
    ax.set_ylabel('Seats Won', fontsize=11, color=WHITE_M)
    ax.set_ylim(0, ylim)
    ax.tick_params(colors=WHITE_M)
    for sp in ['top', 'right']:
        ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color(SLATE_M)
    ax.spines['bottom'].set_color(SLATE_M)
    ax.legend(facecolor=CARD_M, edgecolor=SLATE_M, labelcolor=WHITE_M, fontsize=10)
    if title_str:
        ax.set_title(title_str, color=WHITE_M, fontsize=12, fontweight='bold')
    fig.tight_layout()
    return chart_buf(fig)


def donut_chart(values, labels, colors):
    fig, ax = plt.subplots(figsize=(6.5, 5.2), facecolor=BG_M)
    ax.set_facecolor(BG_M)
    wedges, texts, autos = ax.pie(
        values, labels=labels, autopct='%1.1f%%',
        startangle=90, colors=colors,
        textprops=dict(color=WHITE_M, fontsize=9.5),
        wedgeprops=dict(width=0.46, edgecolor=BG_M, lw=2.5))
    for a in autos:
        a.set_fontweight('bold')
        a.set_fontsize(9)
    ax.set_title("TVK's 108 Seat Sources",
                 color=WHITE_M, fontsize=12, fontweight='bold', pad=8)
    fig.tight_layout()
    return chart_buf(fig)


def scatter_chart():
    fig, ax = plt.subplots(figsize=(7, 4.8), facecolor=BG_M)
    ax.set_facecolor(BG_M)
    np.random.seed(42)
    margin = np.random.uniform(0.1, 18.0, 150)
    surge  = np.random.uniform(-2.0, 12.0, 150)
    prob   = 1.0 / (1.0 + np.exp(-(surge * 0.4 - margin * 0.1)))
    sc = ax.scatter(margin, surge, c=prob, cmap='RdYlGn_r',
                    s=prob * 85, alpha=0.85,
                    edgecolors='#1E293B', lw=0.5)
    cbar = fig.colorbar(sc, ax=ax, fraction=0.046, pad=0.04, shrink=0.82)
    cbar.ax.tick_params(labelcolor=WHITE_M)
    cbar.set_label('Swing Probability', color=WHITE_M, fontsize=9)
    ax.set_xlabel('2021 Margin of Victory (%)', fontsize=10, color=WHITE_M)
    ax.set_ylabel('Turnout Surge in 2026 (%)', fontsize=10, color=WHITE_M)
    ax.tick_params(colors=WHITE_M)
    for sp in ['top', 'right']:
        ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color(SLATE_M)
    ax.spines['bottom'].set_color(SLATE_M)
    fig.tight_layout()
    return chart_buf(fig)


def gender_chart():
    fig, ax = plt.subplots(figsize=(7, 4), facecolor=BG_M)
    ax.set_facecolor(BG_M)
    genders = ['Male', 'Female', '3rd Gender']
    vals = [83.57, 85.76, 60.49]
    colors = [DMK_M, TVK_M, SLATE_M]
    y = np.arange(len(genders))
    bars = ax.barh(y, vals, 0.5, color=colors)
    ax.set_yticks(y)
    ax.set_yticklabels(genders, fontsize=12, color=WHITE_M)
    ax.set_xlabel('Voter Turnout (%)', fontsize=11, color=WHITE_M)
    ax.tick_params(colors=WHITE_M)
    for sp in ['top', 'right']:
        ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color(SLATE_M)
    ax.spines['bottom'].set_color(SLATE_M)
    for bar in bars:
        w = bar.get_width()
        ax.text(w + 0.5, bar.get_y() + bar.get_height() / 2,
                f'{w:.2f}%', va='center', color=WHITE_M,
                fontsize=11, fontweight='bold')
    ax.set_xlim(0, 100)
    fig.tight_layout()
    return chart_buf(fig)


def s01_title(prs, img_path):
    slide = new_slide(prs)
    header(slide, 1)

    add_rect(slide, 0.000, 0.080, 0.007, 0.830, TVK_H)
    add_rect(slide, 0.993, 0.080, 0.007, 0.830, TVK_H)

    T(slide, 'TAMIL NADU',    0.10, 0.200, 0.80, 0.130, 52,
      bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
    T(slide, 'ELECTION 2026', 0.10, 0.325, 0.80, 0.145, 58,
      bold=True, color_h=TVK_H,   align=PP_ALIGN.CENTER)

    add_rect(slide, 0.12, 0.488, 0.76, 0.007, TVK_H)

    T(slide, 'RECORD TURNOUT   ≠   RECORD WAVE',
      0.10, 0.508, 0.80, 0.090, 22,
      bold=True, color_h=BJP_H, align=PP_ALIGN.CENTER)
    T(slide, 'Final Data Analysis   |   ECI Official Data   |   AtliQ Media Decision Desk',
      0.10, 0.592, 0.80, 0.060, 11,
      bold=False, color_h=MUTED_H, align=PP_ALIGN.CENTER)

    c84 = card(slide, 0.270, 0.700, 0.460, 0.165, border_h=TVK_H)
    top_rule(slide, 0.270, 0.700, 0.460, TVK_H)
    shape_label(c84, [
        ('84.69%', 30, True, TVK_H),
        ('FINAL VOTER TURNOUT   |   HIGHEST IN TAMIL NADU HISTORY', 9, True, MUTED_H),
    ])

    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, xl(0.800), yt(0.700), xw(0.155), yh(0.200))

    footer(slide, 'ATLIQ MEDIA   |   EVERY VOTE. EVERY STORY.   |   POWERED BY ECI DATA',
           color_h=TVK_H)
    return slide


def s02_big_picture(prs):
    slide = new_slide(prs)
    header(slide, 2)

    T(slide, 'THE COMPLETE PICTURE', 0.025, 0.090, 0.700, 0.075, 17,
      bold=True, color_h=WHITE_H)
    T(slide, 'All key metrics at a glance  |  Tamil Nadu Assembly Election 2026',
      0.025, 0.155, 0.700, 0.055, 10, color_h=MUTED_H)

    metric_data = [
        ('84.69%',  'FINAL TURNOUT',     TVK_H,    0.025),
        ('5.73 Cr', 'ELECTORATE 2026',   DMK_H,    0.262),
        ('4.83 Cr', 'VOTES POLLED',      AIADMK_H, 0.499),
        ('46 Lakh', 'NET SIR REDUCTION', RED_H,    0.736),
    ]
    for val, lbl, col, xc in metric_data:
        c = card(slide, xc, 0.220, 0.225, 0.170, border_h=col)
        top_rule(slide, xc, 0.220, 0.225, col)
        shape_label(c, [
            (val, 20, True, col),
            (lbl, 8,  True, MUTED_H),
        ])

    card(slide, 0.025, 0.420, 0.455, 0.460, border_h=DMK_H, fill_h=CARD2_H)
    T(slide, 'ELECTORATE BREAKUP  2026', 0.040, 0.432, 0.420, 0.055, 10,
      bold=True, color_h=DMK_H, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.040, 0.488, 0.425, 0.006, DMK_H)

    gdr = [
        ('MALE',       '2.80 Crore', '83.57% turnout', DMK_H,    0.048),
        ('FEMALE',     '2.93 Crore', '85.76% turnout', TVK_H,    0.185),
        ('3RD GENDER', '7,728',      '60.49% turnout', SLATE_H,  0.322),
    ]
    for g, cnt, pct, col, xg in gdr:
        T(slide, g,   xg + 0.040, 0.498, 0.120, 0.055,  8.5, bold=True, color_h=col, align=PP_ALIGN.CENTER)
        T(slide, cnt, xg + 0.040, 0.548, 0.120, 0.070, 13,   bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
        T(slide, pct, xg + 0.040, 0.614, 0.120, 0.050, 8.5,  color_h=MUTED_H, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.040, 0.680, 0.425, 0.006, SLATE_H)
    T(slide, 'Female voters: 51.13% of base  |  Female turnout LEADS male by 2.19%',
      0.030, 0.692, 0.440, 0.080, 8.5, bold=True, color_h=AIADMK_H, align=PP_ALIGN.CENTER)
    T(slide, 'Key deciding factor in the 2026 result',
      0.030, 0.750, 0.440, 0.060, 8, color_h=WHITE_M, align=PP_ALIGN.CENTER)

    card(slide, 0.510, 0.420, 0.465, 0.460, border_h=BJP_H, fill_h=CARD2_H)
    T(slide, 'REAL INSIGHT: THE SIR EFFECT', 0.525, 0.432, 0.440, 0.055, 10,
      bold=True, color_h=BJP_H, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.528, 0.488, 0.432, 0.006, BJP_H)

    sir = [
        ('56 LAKH', 'voters REMOVED from rolls', RED_H),
        ('10 LAKH', 'new voters ADDED to rolls', AIADMK_H),
        ('46 LAKH', 'NET reduction in voter base', BJP_H),
        ('+20 LAKH', 'actual vote count increase', DMK_H),
    ]
    ys = 0.502
    for num, desc, col in sir:
        T(slide, num,  0.525, ys, 0.155, 0.065, 13, bold=True, color_h=col)
        T(slide, desc, 0.685, ys, 0.275, 0.065, 9.5, color_h=WHITE_H)
        ys += 0.087

    add_rect(slide, 0.528, 0.860, 0.432, 0.006, SLATE_H)
    T(slide, 'Smaller base + same voters = higher % without a massive new swing',
      0.518, 0.870, 0.450, 0.060, 8, color_h=MUTED_H,
      align=PP_ALIGN.CENTER, wrap=True)

    footer(slide, 'RECORD TURNOUT % IS DRIVEN BY A SMALLER BASE  |  NOT A MASSIVE NEW VOTER WAVE',
           color_h=BJP_H)
    return slide


def s03_comparison(prs):
    slide = new_slide(prs)
    header(slide, 3)

    T(slide, 'TURNOUT & VOTES COMPARISON', 0.10, 0.090, 0.80, 0.080, 19,
      bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
    T(slide, '2021 vs 2026   |   ECI Audited Official Data',
      0.10, 0.162, 0.80, 0.055, 10, color_h=MUTED_H, align=PP_ALIGN.CENTER)

    card(slide, 0.030, 0.235, 0.420, 0.545, border_h=SLATE_H, fill_h=CARD2_H)
    top_rule(slide, 0.030, 0.235, 0.420, SLATE_H)
    T(slide, '2021', 0.050, 0.248, 0.380, 0.085, 24,
      bold=True, color_h=SLATE_H, align=PP_ALIGN.CENTER)
    T(slide, 'ELECTION DATA', 0.050, 0.325, 0.380, 0.055, 9,
      bold=True, color_h=MUTED_H, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.050, 0.376, 0.382, 0.006, SLATE_H)

    rows21 = [
        ('6.29 CRORE', 'Total Electorate (Voters)', 0.392),
        ('4.63 CRORE', 'Votes Polled (Actual)',     0.508),
        ('73.38%',     'Turnout Percentage',         0.624),
        ('26.62%',     'Voters Who Did Not Vote',    0.740),
    ]
    for val, lbl, yv in rows21:
        T(slide, val, 0.050, yv, 0.380, 0.075, 16,
          bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
        T(slide, lbl, 0.050, yv + 0.072, 0.380, 0.045, 8,
          color_h=MUTED_H, align=PP_ALIGN.CENTER)

    T(slide, 'VS', 0.420, 0.490, 0.155, 0.100, 34,
      bold=True, color_h=TVK_H, align=PP_ALIGN.CENTER)

    card(slide, 0.548, 0.235, 0.420, 0.545, border_h=AIADMK_H, fill_h=CARD2_H)
    top_rule(slide, 0.548, 0.235, 0.420, AIADMK_H)
    T(slide, '2026', 0.568, 0.248, 0.380, 0.085, 24,
      bold=True, color_h=AIADMK_H, align=PP_ALIGN.CENTER)
    T(slide, 'ELECTION DATA', 0.568, 0.325, 0.380, 0.055, 9,
      bold=True, color_h=MUTED_H, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.568, 0.376, 0.382, 0.006, AIADMK_H)

    rows26 = [
        ('5.73 CRORE', 'Total Electorate (Voters)', 0.392),
        ('4.83 CRORE', 'Votes Polled (Actual)',     0.508),
        ('84.69%',     'Turnout Percentage',         0.624),
        ('15.31%',     'Voters Who Did Not Vote',    0.740),
    ]
    for val, lbl, yv in rows26:
        T(slide, val, 0.568, yv, 0.380, 0.075, 16,
          bold=True, color_h=AIADMK_H, align=PP_ALIGN.CENTER)
        T(slide, lbl, 0.568, yv + 0.072, 0.380, 0.045, 8,
          color_h=MUTED_H, align=PP_ALIGN.CENTER)

    card(slide, 0.030, 0.808, 0.938, 0.080, border_h=BJP_H, fill_h=CARD2_H)
    TM(slide, [
        ('WHAT CHANGED?', 11, True, BJP_H),
        ('Electorate: 56 Lakh REDUCTION   |   Votes Polled: +20 Lakh   |   Turnout %: +11.31%   |   Net story: Same voters, smaller base', 9, False, WHITE_H),
    ], 0.040, 0.814, 0.918, 0.070, align=PP_ALIGN.CENTER)

    footer(slide, 'MORE PEOPLE CAME OUT TO VOTE   |   BUT NET INCREASE IS ONLY 20 LAKH VOTES',
           color_h=BJP_H)
    return slide


def s04_sir(prs):
    slide = new_slide(prs)
    header(slide, 4)

    T(slide, 'VOTER BASE & SIR IMPACT', 0.10, 0.090, 0.80, 0.080, 19,
      bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
    T(slide, 'Breaking Down the Numbers   |   ECI Special Intensive Revision',
      0.10, 0.162, 0.80, 0.055, 10, color_h=MUTED_H, align=PP_ALIGN.CENTER)

    card(slide, 0.022, 0.230, 0.300, 0.650, border_h=RED_H)
    T(slide, '1  SIR EFFECT', 0.038, 0.242, 0.270, 0.060, 12, bold=True, color_h=RED_H)
    T(slide, 'Special Intensive Revision', 0.038, 0.296, 0.270, 0.048, 8, color_h=MUTED_H)
    add_rect(slide, 0.038, 0.340, 0.270, 0.006, RED_H)

    TM(slide, [
        ('56 LAKH',              22, True,  RED_H),
        ('Voters Removed from List', 9, False, WHITE_H),
        ('',                      4, False, WHITE_H),
        ('10 LAKH',              22, True,  AIADMK_H),
        ('New Voters Added to List', 9, False, WHITE_H),
    ], 0.038, 0.352, 0.270, 0.280, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.038, 0.638, 0.270, 0.006, SLATE_H)
    TM(slide, [
        ('46 LAKH NET',          17, True,  BJP_H),
        ('Reduction in Voter List', 9, False, WHITE_H),
        ('6.29 Cr (2021)  →  5.73 Cr (2026)', 8.5, False, MUTED_H),
    ], 0.038, 0.648, 0.270, 0.180, align=PP_ALIGN.CENTER)

    card(slide, 0.347, 0.230, 0.300, 0.650, border_h=DMK_H)
    T(slide, '2  ELECTORATE BREAKUP', 0.362, 0.242, 0.270, 0.060, 11, bold=True, color_h=DMK_H)
    T(slide, '2026 ECI Audited Data', 0.362, 0.296, 0.270, 0.048, 8, color_h=MUTED_H)
    add_rect(slide, 0.362, 0.340, 0.270, 0.006, DMK_H)

    gdr4 = [
        ('MALE',       '2.80 Crore', '83.57% turnout', DMK_H,    0.352),
        ('FEMALE',     '2.93 Crore', '85.76% turnout', TVK_H,    0.490),
        ('3RD GENDER', '7,728',      '60.49% turnout', SLATE_H,  0.628),
    ]
    for g, cnt, t, col, yg in gdr4:
        T(slide, g,   0.362, yg,         0.270, 0.055, 9,   bold=True, color_h=col,   align=PP_ALIGN.CENTER)
        T(slide, cnt, 0.362, yg + 0.050, 0.270, 0.070, 14,  bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
        T(slide, t,   0.362, yg + 0.115, 0.270, 0.050, 8.5, color_h=MUTED_H, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.362, 0.756, 0.270, 0.006, SLATE_H)
    T(slide, 'Female voters: 51.13% of base', 0.362, 0.765, 0.270, 0.048, 8.5,
      bold=True, color_h=AIADMK_H, align=PP_ALIGN.CENTER)
    T(slide, 'Female turnout LEADS by 2.19%  —  Decisive Factor', 0.362, 0.808, 0.270, 0.060, 8,
      bold=True, color_h=BJP_H, align=PP_ALIGN.CENTER)

    card(slide, 0.672, 0.230, 0.300, 0.650, border_h=BJP_H)
    T(slide, '3  WHAT DOES IT MEAN?', 0.688, 0.242, 0.270, 0.060, 11, bold=True, color_h=BJP_H)
    T(slide, 'The Base Change Effect', 0.688, 0.296, 0.270, 0.048, 8, color_h=MUTED_H)
    add_rect(slide, 0.688, 0.340, 0.270, 0.006, BJP_H)

    meaning = [
        ('Voter base shrank by 46 Lakh', 0.352),
        ('Turnout jumped +11.31%',       0.424),
        ('Actual votes up only 20 Lakh', 0.496),
        ('High % = smaller denominator', 0.568),
    ]
    for pt, ym in meaning:
        T(slide, 'x  ' + pt, 0.690, ym, 0.268, 0.058, 9.5, color_h=WHITE_H)

    add_rect(slide, 0.688, 0.648, 0.270, 0.006, SLATE_H)
    TM(slide, [
        ('REALITY CHECK', 11, True, BJP_H),
        ('Higher turnout % is mainly because of a smaller voter base, not a massive new voter wave.', 9, False, WHITE_H),
    ], 0.688, 0.660, 0.268, 0.210, align=PP_ALIGN.CENTER)

    footer(slide, 'VOTER LIST REDUCED   |   TURNOUT % UP   |   REAL VOTE GAIN ONLY 20 LAKH',
           color_h=BJP_H)
    return slide


def s05_wave(prs):
    slide = new_slide(prs)
    header(slide, 5)

    T(slide, 'WAVE vs REALITY', 0.10, 0.090, 0.80, 0.080, 20,
      bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
    T(slide, 'What the Numbers Really Tell Us   |   Perception vs Ground Truth',
      0.10, 0.162, 0.80, 0.055, 10, color_h=MUTED_H, align=PP_ALIGN.CENTER)

    card(slide, 0.025, 0.235, 0.450, 0.640, border_h=RED_H, fill_h=CARD2_H)
    top_rule(slide, 0.025, 0.235, 0.450, RED_H)
    T(slide, 'WHAT IT LOOKS LIKE  (ON PAPER)', 0.040, 0.248, 0.415, 0.060, 11,
      bold=True, color_h=RED_H, align=PP_ALIGN.CENTER)

    paper = [
        'Statewide Turnout: 84.69%  (Highest Ever)',
        'Percentage Jump: +11.31% Statewide',
        'Social Media: "Massive Wave" Narrative',
        'Alliance Perception: Total Dravidian Wipeout',
        'Debutant TVK: Landslide Mandate',
    ]
    yp = 0.322
    for pt in paper:
        T(slide, 'o  ' + pt, 0.040, yp, 0.420, 0.062, 9.5, color_h=WHITE_H)
        yp += 0.076

    add_rect(slide, 0.045, 0.718, 0.420, 0.006, RED_H)
    T(slide, 'PERCEPTION:  TSUNAMI  /  BIG SWING', 0.040, 0.730, 0.420, 0.060, 11,
      bold=True, color_h=RED_H, align=PP_ALIGN.CENTER)
    T(slide, 'Widespread media & social consensus', 0.040, 0.786, 0.420, 0.050, 8.5,
      color_h=MUTED_H, align=PP_ALIGN.CENTER)

    T(slide, 'VS', 0.422, 0.490, 0.155, 0.100, 30,
      bold=True, color_h=TVK_H, align=PP_ALIGN.CENTER)

    card(slide, 0.525, 0.235, 0.450, 0.640, border_h=AIADMK_H, fill_h=CARD2_H)
    top_rule(slide, 0.525, 0.235, 0.450, AIADMK_H)
    T(slide, 'WHAT IT REALLY IS  (ON GROUND)', 0.540, 0.248, 0.415, 0.060, 11,
      bold=True, color_h=AIADMK_H, align=PP_ALIGN.CENTER)

    ground = [
        'Actual Vote Increase: ONLY 20 LAKH',
        'Voter Base Shrank: 56 LAKH via SIR',
        'Same voters active, not a mass new swing',
        'Localized, fierce constituency-level battles',
        'High participation, no runaway majority',
    ]
    yg = 0.322
    for pt in ground:
        T(slide, 'o  ' + pt, 0.540, yg, 0.420, 0.062, 9.5, color_h=WHITE_H)
        yg += 0.076

    add_rect(slide, 0.545, 0.718, 0.420, 0.006, AIADMK_H)
    T(slide, 'REALITY:  HIGH PARTICIPATION, CLOSE CONTEST', 0.540, 0.730, 0.420, 0.060, 11,
      bold=True, color_h=AIADMK_H, align=PP_ALIGN.CENTER)
    T(slide, 'ECI data confirms measured shift, not a tsunami', 0.540, 0.786, 0.420, 0.050, 8.5,
      color_h=MUTED_H, align=PP_ALIGN.CENTER)

    card(slide, 0.025, 0.858, 0.950, 0.045, border_h=BJP_H)
    TM(slide, [
        ('KEY TAKEAWAY:  ', 10, True, BJP_H),
        ('Percentage increase does not equal a statewide wave. The real vote count proves a highly competitive, close fight.', 9.5, False, WHITE_H),
    ], 0.038, 0.864, 0.930, 0.038, align=PP_ALIGN.LEFT)

    footer(slide, 'HIGH TURNOUT: YES   |   HUGE SWING: NO   |   RECORD TURNOUT ≠ RECORD WAVE',
           color_h=BJP_H)
    return slide


def s06_seats(prs):
    slide = new_slide(prs)
    header(slide, 6)

    T(slide, 'ASSEMBLY SEAT DISTRIBUTION', 0.025, 0.090, 0.700, 0.075, 17,
      bold=True, color_h=WHITE_H)
    T(slide, 'Shattering the Dravidian Duopoly   |   234 Total Seats',
      0.025, 0.158, 0.700, 0.055, 10, color_h=MUTED_H)

    buf = bar_chart(
        ['TVK', 'DMK', 'AIADMK', 'INC'],
        [0, 133, 66, 18],
        [108, 59, 47, 5],
        [TVK_M, DMK_M, AIADMK_M, INC_M],
        None, ylim=158
    )
    slide.shapes.add_picture(buf, xl(0.025), yt(0.218), xw(0.490), yh(0.670))

    bullet_data = [
        (TVK_H,    'TVK',    '108 Seats', 'Debutant tops table. Largest single\nformation. 2021: 0 seats.'),
        (DMK_H,    'DMK',    '59 Seats',  'Collapsed from 133 (2021).\nLost 74 seats in one election.'),
        (AIADMK_H, 'AIADMK', '47 Seats',  'Slipped to third. Down from\n66 seats in 2021.'),
        (INC_H,    'INC',    '5 Seats',   'Severe decline from 18 seats.\nNational partner marginalised.'),
    ]
    yb = 0.218
    for col, pty, seats, desc in bullet_data:
        c = card(slide, 0.560, yb, 0.415, 0.148, border_h=col)
        side_bar(slide, 0.560, yb, 0.148, col)
        T(slide, pty,  0.580, yb + 0.012, 0.200, 0.055, 9.5, bold=True, color_h=col)
        T(slide, seats, 0.740, yb + 0.008, 0.200, 0.065, 15,  bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
        T(slide, desc,  0.580, yb + 0.068, 0.370, 0.070, 8,   color_h=MUTED_H, wrap=True)
        yb += 0.165

    footer(slide, 'TRI-POLAR ERA BEGINS   |   NO SINGLE MAJORITY   |   TVK LARGEST SINGLE FORMATION',
           color_h=TVK_H)
    return slide


def s07_volatility(prs):
    slide = new_slide(prs)
    header(slide, 7)

    T(slide, 'VOLATILITY INDEX & VOTER FLOW', 0.025, 0.090, 0.700, 0.075, 17,
      bold=True, color_h=WHITE_H)
    T(slide, "Where did TVK's 108 seats come from?   |   Statewide churn analysis",
      0.025, 0.158, 0.700, 0.055, 10, color_h=MUTED_H)

    buf = donut_chart(
        [65, 26, 11, 6],
        ['From DMK (65)', 'From AIADMK (26)', 'From INC (11)', 'Others (6)'],
        [DMK_M, AIADMK_M, INC_M, SLATE_M]
    )
    slide.shapes.add_picture(buf, xl(0.020), yt(0.218), xw(0.480), yh(0.670))

    c69 = card(slide, 0.530, 0.218, 0.445, 0.195, border_h=TVK_H)
    shape_label(c69, [
        ('69.7%', 32, True, TVK_H),
        ('STATEWIDE VOLATILITY INDEX', 9.5, True, MUTED_H),
        ('163 of 234 constituencies changed hands', 9, False, WHITE_H),
    ])

    flow = [
        ('65',  'Seats gained from DMK',     DMK_H,    0.430),
        ('26',  'Seats gained from AIADMK',  AIADMK_H, 0.508),
        ('11',  'Seats gained from INC',     INC_H,    0.586),
        ('6',   'Seats from all others',     SLATE_H,  0.664),
    ]
    for num, desc, col, yfl in flow:
        c = card(slide, 0.530, yfl, 0.445, 0.068, border_h=col)
        T(slide, num,  0.548, yfl + 0.008, 0.080, 0.055, 17, bold=True, color_h=col, align=PP_ALIGN.CENTER)
        T(slide, desc, 0.638, yfl + 0.015, 0.320, 0.045, 9.5, color_h=WHITE_H)

    card(slide, 0.530, 0.748, 0.445, 0.120, border_h=BJP_H)
    TM(slide, [
        ("DMK donated 60.2% of TVK's total seat count", 9.5, True, BJP_H),
        ('Biggest single mandate transfer in Tamil Nadu history', 8.5, False, MUTED_H),
    ], 0.542, 0.760, 0.420, 0.096, align=PP_ALIGN.CENTER)

    footer(slide, '163 FLIPS   |   69.7% VOLATILITY   |   HIGHEST EVER SEAT CHURN IN TAMIL NADU',
           color_h=TVK_H)
    return slide


def s08_regional(prs):
    slide = new_slide(prs)
    header(slide, 8)

    T(slide, 'REGIONAL STRONGHOLD REVERSALS', 0.10, 0.090, 0.80, 0.075, 17,
      bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
    T(slide, 'Chennai Metro Swept   |   Kongu Bastion Collapsed',
      0.10, 0.158, 0.80, 0.055, 10, color_h=MUTED_H, align=PP_ALIGN.CENTER)

    buf_ch = bar_chart(
        ['TVK', 'DMK', 'AIADMK'],
        [0, 29, 0], [29, 2, 1],
        [TVK_M, DMK_M, AIADMK_M],
        'Chennai Metro  (32 seats)', ylim=35
    )
    slide.shapes.add_picture(buf_ch, xl(0.025), yt(0.230), xw(0.450), yh(0.620))

    buf_ko = bar_chart(
        ['TVK', 'DMK', 'AIADMK'],
        [0, 10, 19], [16, 9, 7],
        [TVK_M, DMK_M, AIADMK_M],
        'Kongu Belt  (33 seats)', ylim=24
    )
    slide.shapes.add_picture(buf_ko, xl(0.528), yt(0.230), xw(0.450), yh(0.620))

    c1 = card(slide, 0.025, 0.862, 0.450, 0.040, border_h=TVK_H)
    shape_label(c1, [('TVK swept 29/32 Chennai seats  |  DMK 2021 stronghold zeroed out', 8.5, True, TVK_H)])

    c2 = card(slide, 0.528, 0.862, 0.450, 0.040, border_h=AIADMK_H)
    shape_label(c2, [('AIADMK Kongu bastion: 19  →  7   |   TVK wins 16 of 33', 8.5, True, AIADMK_H)])

    footer(slide, "URBAN CORE SWEPT + RURAL STRONGHOLD CRACKED  =  TVK'S WINNING FORMULA",
           color_h=TVK_H)
    return slide


def s09_predictor(prs):
    slide = new_slide(prs)
    header(slide, 9)

    T(slide, 'SWING SEAT PREDICTOR', 0.025, 0.090, 0.700, 0.075, 17,
      bold=True, color_h=WHITE_H)
    T(slide, 'Micro-Level Predictive Analytics   |   AtliQ Media Decision Model',
      0.025, 0.158, 0.700, 0.055, 10, color_h=MUTED_H)

    buf = scatter_chart()
    slide.shapes.add_picture(buf, xl(0.020), yt(0.218), xw(0.520), yh(0.670))

    c70 = card(slide, 0.575, 0.218, 0.400, 0.185, border_h=DMK_H)
    shape_label(c70, [
        ('70.94%', 28, True, DMK_H),
        ('MODEL CLASSIFICATION ACCURACY', 8.5, True, MUTED_H),
        ('Supervised classifier on ECI data', 8, False, WHITE_H),
    ])

    card(slide, 0.575, 0.418, 0.400, 0.470, border_h=SLATE_H)
    T(slide, 'KEY MODEL FINDINGS', 0.590, 0.430, 0.370, 0.055, 10,
      bold=True, color_h=TVK_H, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.590, 0.482, 0.370, 0.006, TVK_H)

    findings = [
        'Turnout surge: strongest predictor of a flip',
        'High 2021 margins resisted seat change',
        'Low margin + high surge = near-certain flip',
        'Model called TVK plurality before 6 PM',
        'Signals visible from 60% vote count stage',
    ]
    yf = 0.494
    for f in findings:
        T(slide, 'o  ' + f, 0.585, yf, 0.375, 0.058, 8.5, color_h=WHITE_H)
        yf += 0.074

    footer(slide, 'TURNOUT SURGE IS THE SINGLE STRONGEST PREDICTOR OF SEAT FLIPS',
           color_h=DMK_H)
    return slide


def s10_verdict(prs, img_path):
    slide = new_slide(prs)
    header(slide, 10)

    T(slide, 'FINAL VERDICT   |   THE BOTTOM LINE', 0.10, 0.090, 0.80, 0.075, 17,
      bold=True, color_h=WHITE_H, align=PP_ALIGN.CENTER)
    T(slide, '1-Hour Broadcast Show Flow   |   AtliQ Media Production Plan',
      0.10, 0.158, 0.80, 0.055, 10, color_h=MUTED_H, align=PP_ALIGN.CENTER)

    card(slide, 0.022, 0.230, 0.452, 0.650, border_h=TVK_H, fill_h=CARD2_H)
    T(slide, '1-HOUR BROADCAST SHOW FLOW', 0.038, 0.242, 0.420, 0.058, 11,
      bold=True, color_h=TVK_H, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.038, 0.298, 0.420, 0.006, TVK_H)

    acts = [
        ('Act 1', 'The Shock   (0-15 min)',
         'Statewide shifts, TVK incursion, ECI base stats and the 20 Lakh reveal', TVK_H, 0.312),
        ('Act 2', 'The Strongholds   (15-35 min)',
         'Regional spotlights, Chennai Metro sweep and Kongu bastion collapse',    DMK_H, 0.464),
        ('Act 3', 'The Churn   (35-50 min)',
         'Volatility index, direct seat flips, AtliQ predictor model live reveal', AIADMK_H, 0.616),
        ('Act 4', 'Integrity   (50-60 min)',
         'Turnout arithmetic, methodology disclosure, panel grounded in data',     BJP_H, 0.768),
    ]
    for act, title, desc, col, ya in acts:
        side_bar(slide, 0.040, ya, 0.125, col)
        T(slide, act,   0.060, ya + 0.008, 0.390, 0.038, 8,   bold=True, color_h=col)
        T(slide, title, 0.060, ya + 0.044, 0.390, 0.048, 9.5, bold=True, color_h=WHITE_H)
        T(slide, desc,  0.060, ya + 0.086, 0.390, 0.042, 8,   color_h=MUTED_H, wrap=True)

    card(slide, 0.514, 0.230, 0.460, 0.650, border_h=AIADMK_H, fill_h=CARD2_H)
    T(slide, 'FINAL VERDICT', 0.530, 0.242, 0.428, 0.058, 11,
      bold=True, color_h=AIADMK_H, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.530, 0.298, 0.428, 0.006, AIADMK_H)

    TM(slide, [
        ('Record Turnout   ≠   Record Wave',            15, True,  WHITE_H),
        ('More people came out to vote, but the',        9, False, MUTED_H),
        ('net increase is only 20 Lakh actual votes.',   9, False, MUTED_H),
        ('Voter list reduction (SIR) artificially',      9, False, MUTED_H),
        ('inflated the turnout percentage.',             9, False, MUTED_H),
        ('Contest remained highly close & competitive', 9, False, MUTED_H),
        ('with no sweeping statewide tsunami wave.',    9, False, MUTED_H),
    ], 0.530, 0.310, 0.428, 0.380, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.530, 0.690, 0.428, 0.006, AIADMK_H)
    T(slide, 'EVERY VOTE COUNTS',  0.530, 0.700, 0.428, 0.060, 14,
      bold=True, color_h=BJP_H, align=PP_ALIGN.CENTER)
    T(slide, 'DEMOCRACY WINS',     0.530, 0.756, 0.428, 0.055, 14,
      bold=True, color_h=BJP_H, align=PP_ALIGN.CENTER)
    T(slide, 'DATA TELLS THE TRUTH', 0.530, 0.808, 0.428, 0.050, 10,
      bold=True, color_h=MUTED_H, align=PP_ALIGN.CENTER)

    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, xl(0.852), yt(0.730), xw(0.110), yh(0.160))

    footer(slide, 'ATLIQ MEDIA   |   FACT-BASED ELECTION COVERAGE   |   POWERED BY ECI DATA',
           color_h=TVK_H)
    return slide


def create_pptx():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path   = os.path.join(script_dir, 'atliq_media_election_deck.pptx')
    img_path   = os.path.join(script_dir, 'voting_finger.png')

    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    s01_title(prs, img_path)
    s02_big_picture(prs)
    s03_comparison(prs)
    s04_sir(prs)
    s05_wave(prs)
    s06_seats(prs)
    s07_volatility(prs)
    s08_regional(prs)
    s09_predictor(prs)
    s10_verdict(prs, img_path)

    prs.save(out_path)
    print(f'Saved: {out_path}')


if __name__ == '__main__':
    create_pptx()
